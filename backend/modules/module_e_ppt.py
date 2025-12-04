"""Module E - Programmatic PPTX builder.

This module assembles PNG charts (Module C) and AI-generated texts (Module D)
into a themed PowerPoint presentation using python-pptx only. It exposes a
single entry point: ``build_presentation``.
"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DEFAULT_FALLBACK_TEXT = "Analyse non disponible."
THEME_PRESETS = {
    "corporate": {
        "name": "corporate",
        "background": "FFFFFF",
        "accent": "1F4E79",
        "accent_light": "E8EEF5",
        "text": "1A1A1A",
        "subtitle": "4A4A4A",
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_align": PP_ALIGN.LEFT,
        "text_box_position": "bottom",
    },
    "minimal": {
        "name": "minimal",
        "background": "FFFFFF",
        "accent": "000000",
        "accent_light": "F5F5F5",
        "text": "111111",
        "subtitle": "555555",
        "title_font": "Inter",
        "body_font": "Inter",
        "title_align": PP_ALIGN.CENTER,
        "text_box_position": "side",
    },
    "dark": {
        "name": "dark",
        "background": "1A1C1F",
        "accent": "0CA6E9",
        "accent_light": "2A2D34",
        "text": "FFFFFF",
        "subtitle": "D1D5DB",
        "title_font": "Montserrat",
        "body_font": "Montserrat",
        "title_align": PP_ALIGN.LEFT,
        "text_box_position": "bottom",
    },
    "vibrant": {
        "name": "vibrant",
        "background": "FFFFFF",
        "accent": "FF5E5B",
        "accent_light": "FFE066",
        "text": "1F2933",
        "subtitle": "475569",
        "title_font": "Montserrat",
        "body_font": "Montserrat",
        "title_align": PP_ALIGN.LEFT,
        "text_box_position": "side",
    },
}


def get_theme_config(theme_name: str) -> Tuple[Dict[str, Any], Optional[str]]:
    theme = THEME_PRESETS.get(theme_name.lower())
    warning = None
    if not theme:
        warning = f"Thème '{theme_name}' non reconnu. Utilisation du thème minimal."
        theme = THEME_PRESETS["minimal"]
    return theme, warning


def create_title_slide(prs: Presentation, title: str, theme_cfg: Dict[str, Any], options: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _apply_background(slide, theme_cfg)

    width = prs.slide_width
    top_bar = slide.shapes.add_shape(
        autoshape_type_id=1,  # rectangle
        left=0,
        top=0,
        width=width,
        height=Inches(0.35),
    )
    _fill_shape(top_bar, theme_cfg["accent"])

    title_box = slide.shapes.add_textbox(left=Inches(0.7), top=Inches(1.2), width=width - Inches(1.4), height=Inches(2.0))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.name = theme_cfg["title_font"]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = _rgb(theme_cfg["text"] if theme_cfg["name"] != "dark" else "FFFFFF")
    p.alignment = theme_cfg.get("title_align", PP_ALIGN.LEFT)

    subtitle = options.get("subtitle") or "Rapport généré automatiquement"
    date_str = datetime.now().strftime("%d %B %Y")
    subtitle_box = slide.shapes.add_textbox(left=Inches(0.7), top=Inches(2.7), width=width - Inches(1.4), height=Inches(1.0))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"{subtitle} — {date_str}"
    subtitle_frame.paragraphs[0].font.size = Pt(18)
    subtitle_frame.paragraphs[0].font.name = theme_cfg["body_font"]
    subtitle_frame.paragraphs[0].font.color.rgb = _rgb(theme_cfg["subtitle"])

    logo_path = options.get("logo_path")
    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(
            logo_path,
            left=width - Inches(2.0),
            top=Inches(0.2),
            width=Inches(1.5),
        )


def create_graph_slide(
    prs: Presentation,
    plot_meta: Dict[str, Any],
    text_meta: Dict[str, Any],
    theme_cfg: Dict[str, Any],
    errors: List[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_background(slide, theme_cfg)
    width, height = prs.slide_width, prs.slide_height

    title = text_meta.get("title") or f"Analyse de {plot_meta.get('column') or 'colonne'}"
    title_box = slide.shapes.add_textbox(left=Inches(0.6), top=Inches(0.4), width=width - Inches(1.2), height=Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = theme_cfg["title_font"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _rgb(theme_cfg["text"] if theme_cfg["name"] != "dark" else "FFFFFF")

    image_path = plot_meta.get("file_path")
    image_exists = image_path and Path(image_path).exists()
    image_height = height - Inches(3.2)
    image_width = width - Inches(1.5)
    image_left = Inches(0.75)
    image_top = Inches(1.3)

    if image_exists:
        _add_image_within_bounds(slide, image_path, image_left, image_top, image_width, image_height)
    else:
        errors.append(f"Image introuvable pour {plot_meta.get('column')} ({plot_meta.get('graph_type')}).")
        placeholder = slide.shapes.add_shape(
            1,
            left=image_left,
            top=image_top,
            width=image_width,
            height=image_height,
        )
        _fill_shape(placeholder, theme_cfg["accent_light"])
        text_frame = placeholder.text_frame
        text_frame.text = "Graphique non disponible"
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    text = text_meta.get("text") or DEFAULT_FALLBACK_TEXT
    text_box_height = Inches(1.8)
    text_box_top = height - text_box_height - Inches(0.5)

    if theme_cfg.get("text_box_position") == "side":
        text_box = slide.shapes.add_textbox(
            left=width - Inches(3.6),
            top=Inches(1.3),
            width=Inches(2.8),
            height=height - Inches(2.0),
        )
    else:
        text_box = slide.shapes.add_textbox(
            left=Inches(0.75),
            top=text_box_top,
            width=width - Inches(1.5),
            height=text_box_height,
        )

    frame = text_box.text_frame
    frame.word_wrap = True
    frame.margin_bottom = 0
    frame.margin_top = 0
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    paragraph = frame.paragraphs[0]
    paragraph.text = _truncate_text(text, 650)
    paragraph.font.size = Pt(16)
    paragraph.font.name = theme_cfg["body_font"]
    paragraph.font.color.rgb = _rgb(theme_cfg["text"] if theme_cfg["name"] != "dark" else "FFFFFF")


def create_data_overview_slide(prs: Presentation, diagnostic: Dict[str, Any], theme_cfg: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_background(slide, theme_cfg)
    width = prs.slide_width

    title_box = slide.shapes.add_textbox(left=Inches(0.7), top=Inches(0.6), width=width - Inches(1.4), height=Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Aperçu du dataset"
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = theme_cfg["title_font"]
    title_frame.paragraphs[0].font.color.rgb = _rgb(theme_cfg["text"] if theme_cfg["name"] != "dark" else "FFFFFF")

    bullets_box = slide.shapes.add_textbox(left=Inches(0.9), top=Inches(1.7), width=width - Inches(1.8), height=Inches(4.5))
    bullets_frame = bullets_box.text_frame
    bullets_frame.clear()

    stats = [
        f"Lignes : {diagnostic.get('num_rows', 'n/a')}",
        f"Colonnes : {diagnostic.get('num_cols', 'n/a')}",
    ]
    columns_info = diagnostic.get("columns", {})
    if columns_info:
        highlighted = list(columns_info.keys())[:5]
        stats.append(f"Colonnes clés : {', '.join(highlighted)}")

    for stat in stats:
        p = bullets_frame.add_paragraph()
        p.text = stat
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = theme_cfg["body_font"]
        p.font.color.rgb = _rgb(theme_cfg["text"] if theme_cfg["name"] != "dark" else "FFFFFF")


def create_conclusion_slide(prs: Presentation, conclusion_text: str, theme_cfg: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_background(slide, theme_cfg)
    width = prs.slide_width

    title_box = slide.shapes.add_textbox(left=Inches(0.7), top=Inches(0.6), width=width - Inches(1.4), height=Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.text = "Synthèse et conclusion"
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = theme_cfg["title_font"]
    title_frame.paragraphs[0].font.color.rgb = _rgb(theme_cfg["text"] if theme_cfg["name"] != "dark" else "FFFFFF")

    text_box = slide.shapes.add_textbox(left=Inches(0.9), top=Inches(1.8), width=width - Inches(1.8), height=Inches(4.5))
    frame = text_box.text_frame
    frame.word_wrap = True
    frame.paragraphs[0].text = conclusion_text or DEFAULT_FALLBACK_TEXT
    frame.paragraphs[0].font.size = Pt(20)
    frame.paragraphs[0].font.name = theme_cfg["body_font"]
    frame.paragraphs[0].font.color.rgb = _rgb(theme_cfg["text"] if theme_cfg["name"] != "dark" else "FFFFFF")


def apply_footer_and_brand(prs: Presentation, options: Dict[str, Any], theme_cfg: Dict[str, Any]) -> None:
    footer_text = options.get("footer_text")
    if not footer_text:
        return
    for slide in prs.slides:
        footer_box = slide.shapes.add_textbox(
            left=Inches(0.6),
            top=prs.slide_height - Inches(0.5),
            width=prs.slide_width - Inches(1.2),
            height=Inches(0.3),
        )
        frame = footer_box.text_frame
        frame.text = footer_text
        frame.paragraphs[0].font.size = Pt(12)
        frame.paragraphs[0].font.name = theme_cfg["body_font"]
        frame.paragraphs[0].font.color.rgb = _rgb(theme_cfg["subtitle"])
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def build_presentation(
    title: str,
    plots: List[Dict[str, Any]],
    texts: Dict[str, Any],
    output_path: str,
    theme: str = "corporate",
    options: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    options = options or {}
    theme_cfg, warning = get_theme_config(theme)
    errors: List[str] = []
    if warning:
        errors.append(warning)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    diagnostic = options.get("diagnostic")
    text_entries = texts.get("analyses", []) if isinstance(texts, dict) else []
    text_map = {
        (entry.get("column"), entry.get("graph_type")): entry
        for entry in text_entries
        if isinstance(entry, dict)
    }

    try:
        create_title_slide(prs, title, theme_cfg, options)
    except Exception as exc:  # pragma: no cover
        errors.append(f"Erreur lors de la création de la slide titre: {exc}")

    if diagnostic:
        try:
            create_data_overview_slide(prs, diagnostic, theme_cfg)
        except Exception as exc:  # pragma: no cover
            errors.append(f"Impossible de créer l'aperçu des données: {exc}")

    ordered_plots = _order_plots(plots, options.get("plots_order"))
    for plot in ordered_plots:
        key = (plot.get("column"), plot.get("graph_type"))
        text_meta = text_map.get(key, {"text": DEFAULT_FALLBACK_TEXT, "title": plot.get("column")})
        try:
            create_graph_slide(prs, plot, text_meta, theme_cfg, errors)
        except Exception as exc:  # pragma: no cover
            errors.append(f"Impossible de créer la slide pour {plot.get('column')}: {exc}")

    conclusion_text = texts.get("conclusion") if isinstance(texts, dict) else None
    try:
        create_conclusion_slide(prs, conclusion_text or DEFAULT_FALLBACK_TEXT, theme_cfg)
    except Exception as exc:  # pragma: no cover
        errors.append(f"Impossible de créer la conclusion: {exc}")

    try:
        apply_footer_and_brand(prs, options, theme_cfg)
    except Exception as exc:  # pragma: no cover
        errors.append(f"Impossible d'appliquer le footer: {exc}")

    try:
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_file))
    except Exception as exc:
        errors.append(f"Erreur lors de l'enregistrement du PPTX: {exc}")

    return {
        "pptx_path": str(output_path),
        "slides": len(prs.slides),
        "errors": errors,
    }


# ----------------------------- helper functions ---------------------------- #

def _apply_background(slide, theme_cfg: Dict[str, Any]) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(theme_cfg["background"])


def _fill_shape(shape, hex_color: str) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)
    line = shape.line
    line.color.rgb = _rgb(hex_color)


def _add_image_within_bounds(slide, image_path: str, left, top, max_width, max_height) -> None:
    left = int(left)
    top = int(top)
    max_width = int(max_width)
    max_height = int(max_height)
    image = slide.shapes.add_picture(image_path, left=left, top=top)
    width_ratio = max_width / image.width
    height_ratio = max_height / image.height
    scale = min(width_ratio, height_ratio, 1)
    image.width = int(image.width * scale)
    image.height = int(image.height * scale)
    image.left = left + (max_width - image.width) // 2
    image.top = top + (max_height - image.height) // 2


def _truncate_text(text: str, limit: int) -> str:
    if not text:
        return DEFAULT_FALLBACK_TEXT
    if len(text) <= limit:
        return text
    return text[: limit - 3] + "..."


def _rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _order_plots(plots: List[Dict[str, Any]], plots_order: Optional[List[str]]) -> List[Dict[str, Any]]:
    if not plots_order:
        return plots
    priority = {value: index for index, value in enumerate(plots_order)}
    return sorted(plots, key=lambda p: priority.get(p.get("column"), len(priority)))
